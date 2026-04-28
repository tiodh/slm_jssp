from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Theme: light/white, clean, sleek
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xFA, 0xFA, 0xFA)
INK = RGBColor(0x1F, 0x24, 0x2E)          # near-black title
MUTED = RGBColor(0x5B, 0x63, 0x72)        # body grey
ACCENT = RGBColor(0x2E, 0x6F, 0xDB)       # subtle blue
LINE = RGBColor(0xE5, 0xE7, 0xEB)         # hairline divider

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def set_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=16, color=INK, line_spacing=1.25, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "•  " + item
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_divider(slide, left, top, width):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = LINE


def add_accent_bar(slide, left, top, width=Inches(0.5), height=Pt(3)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT


def add_header(slide, eyebrow, title):
    add_accent_bar(slide, Inches(0.6), Inches(0.55))
    add_text(slide, Inches(0.6), Inches(0.62), Inches(12), Inches(0.3),
             eyebrow, size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.6), Inches(0.88), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=INK)
    add_divider(slide, Inches(0.6), Inches(1.55), Inches(12.15))


def add_footer(slide, page, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             "Starjob · JSSP with LLMs", size=10, color=MUTED)
    add_text(slide, Inches(10.7), Inches(7.05), Inches(2), Inches(0.3),
             f"{page} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# =========================================================================
# Slide 1 — Title
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_accent_bar(s, Inches(0.9), Inches(2.6), width=Inches(1.0), height=Pt(4))
add_text(s, Inches(0.9), Inches(2.7), Inches(12), Inches(0.4),
         "FINE-TUNING LARGE LANGUAGE MODELS", size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.9), Inches(3.1), Inches(12), Inches(1.6),
         "Solving the Job-Shop Scheduling Problem",
         size=44, bold=True, color=INK)
add_text(s, Inches(0.9), Inches(4.0), Inches(12), Inches(0.6),
         "A comparative study across four 7–8B instruction-tuned models",
         size=20, color=MUTED)
add_divider(s, Inches(0.9), Inches(5.0), Inches(11.5))
add_text(s, Inches(0.9), Inches(5.15), Inches(12), Inches(0.4),
         "Starjob Project  ·  LoRA Fine-Tuning  ·  2026",
         size=13, color=MUTED)

# =========================================================================
# Slide 2 — Introduction
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header(s, "01  ·  INTRODUCTION", "Why scheduling meets language models")

add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.4),
         "The Job-Shop Scheduling Problem (JSSP)",
         size=20, bold=True, color=INK)
add_bullets(s, Inches(0.6), Inches(2.4), Inches(12), Inches(2),
            [
                "Classical NP-hard combinatorial problem: assign n jobs across m machines to minimize makespan.",
                "Exact solvers (CP-SAT, MILP) scale poorly; heuristics (dispatching rules, GA) trade optimality for speed.",
                "A natural testbed for reasoning-heavy, structured-output generation with language models.",
            ], size=15, color=MUTED)

add_text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
         "Our motivation",
         size=20, bold=True, color=INK)
add_bullets(s, Inches(0.6), Inches(5.0), Inches(12), Inches(2),
            [
                "Can modern 7–8B LLMs learn to produce feasible, near-optimal schedules from natural-language specs?",
                "How do different model families (Llama, Ministral, Qwen2, Granite) compare under identical training?",
                "Deliver a fair, reproducible benchmark with the Starjob dataset (130k JSSP instances).",
            ], size=15, color=MUTED)

add_footer(s, 2, 6)

# =========================================================================
# Slide 3 — Problem & Dataset
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header(s, "02  ·  PROBLEM FORMULATION", "Starjob: JSSP as instruction-following")

# Left card — Input
card1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.6), Inches(1.9), Inches(6.0), Inches(4.6))
card1.fill.solid()
card1.fill.fore_color.rgb = OFFWHITE
card1.line.color.rgb = LINE
card1.line.width = Pt(0.75)

add_text(s, Inches(0.9), Inches(2.1), Inches(5.5), Inches(0.4),
         "Input (natural language)", size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(0.9), Inches(2.55), Inches(5.5), Inches(3.8),
            [
                "Number of jobs n and machines m.",
                "For each job: ordered sequence of (machine, processing-time) operations.",
                "Alpaca-style instruction + input prompt.",
            ], size=14, color=INK)

# Right card — Output
card2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(6.9), Inches(1.9), Inches(6.0), Inches(4.6))
card2.fill.solid()
card2.fill.fore_color.rgb = OFFWHITE
card2.line.color.rgb = LINE
card2.line.width = Pt(0.75)

add_text(s, Inches(7.2), Inches(2.1), Inches(5.5), Inches(0.4),
         "Output (structured schedule)", size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(7.2), Inches(2.55), Inches(5.5), Inches(3.8),
            [
                "Per-operation start and end times.",
                "Final makespan (lower is better).",
                "Must respect job precedence & machine exclusivity.",
            ], size=14, color=INK)

# Dataset strip
add_text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.4),
         "Dataset — Starjob · 130k instances · sizes from 2×2 up to 20×15 · 2% held-out eval split",
         size=13, bold=True, color=MUTED)

add_footer(s, 3, 6)

# =========================================================================
# Slide 4 — Methodology: Models & Training
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header(s, "03  ·  METHODOLOGY", "Fair comparison across four model families")

add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.4),
         "Models (all instruction-tuned, 7–8B parameters)",
         size=18, bold=True, color=INK)
add_bullets(s, Inches(0.6), Inches(2.4), Inches(12), Inches(2),
            [
                "Llama-3.1-8B-Instruct  —  Meta",
                "Ministral-8B-Instruct  —  Mistral AI",
                "Qwen2-7B-Instruct  —  Alibaba",
                "Granite-3.2-8B-Instruct  —  IBM",
            ], size=15, color=MUTED)

add_text(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.4),
         "Identical training recipe",
         size=18, bold=True, color=INK)
add_bullets(s, Inches(0.6), Inches(5.2), Inches(12), Inches(2),
            [
                "LoRA: r = 32, α = 32, dropout = 0.0, target = all attention + MLP projections.",
                "4-bit NF4 quantization · bf16 compute · max seq length 8192 · Unsloth runtime on RTX 4090.",
                "Batch 1 × grad-accum 8 · AdamW-8bit · lr 2e-4, linear schedule · 1 epoch · seed 42.",
            ], size=15, color=MUTED)

add_footer(s, 4, 6)

# =========================================================================
# Slide 5 — Methodology: Pipeline
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_header(s, "04  ·  PIPELINE", "From instance to evaluated makespan")

steps = [
    ("Starjob\ninstance", "JSON → Alpaca prompt"),
    ("LoRA\nfine-tune", "SFTTrainer (TRL)\n+ Unsloth"),
    ("Generate\nschedule", "Greedy decode\nmax 4096 tokens"),
    ("Parse &\nvalidate", "Regex extract\nfeasibility check"),
    ("Compute\nmetrics", "Gap, exact-match,\nfeasibility"),
]

x0 = Inches(0.6)
y0 = Inches(2.8)
w = Inches(2.25)
h = Inches(1.8)
gap = Inches(0.18)

for i, (title, sub) in enumerate(steps):
    left = x0 + (w + gap) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y0, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = OFFWHITE
    box.line.color.rgb = LINE
    box.line.width = Pt(0.75)

    # step number
    add_text(s, left + Inches(0.2), y0 + Inches(0.15), Inches(1), Inches(0.3),
             f"STEP {i+1}", size=10, bold=True, color=ACCENT)
    # title
    add_text(s, left + Inches(0.2), y0 + Inches(0.45), w - Inches(0.4), Inches(0.8),
             title, size=15, bold=True, color=INK)
    # subtitle
    add_text(s, left + Inches(0.2), y0 + Inches(1.15), w - Inches(0.4), Inches(0.7),
             sub, size=11, color=MUTED)

    # arrow between boxes
    if i < len(steps) - 1:
        ax = left + w + Inches(0.01)
        ay = y0 + h / 2 - Pt(5)
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, gap - Inches(0.02), Pt(10))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

# Caption
add_text(s, Inches(0.6), Inches(5.1), Inches(12), Inches(0.5),
         "Each model runs through the same pipeline with the same hyper-parameters — "
         "differences in the final metrics are attributable to the base model, not the training setup.",
         size=13, color=MUTED)

# Evaluation metrics strip
add_text(s, Inches(0.6), Inches(5.9), Inches(12), Inches(0.4),
         "Evaluation metrics", size=16, bold=True, color=INK)
add_bullets(s, Inches(0.6), Inches(6.3), Inches(12), Inches(1),
            [
                "Optimality gap %  ·  Exact-match rate  ·  Feasibility rate  ·  MAE  ·  Breakdown by problem size (small / medium / large).",
            ], size=12, color=MUTED)

add_footer(s, 5, 6)

# =========================================================================
# Slide 6 — Closing divider
# =========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_accent_bar(s, Inches(0.9), Inches(3.2), width=Inches(1.0), height=Pt(4))
add_text(s, Inches(0.9), Inches(3.3), Inches(12), Inches(0.4),
         "NEXT", size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.9), Inches(3.7), Inches(12), Inches(1.2),
         "Results & Discussion",
         size=44, bold=True, color=INK)
add_text(s, Inches(0.9), Inches(4.7), Inches(12), Inches(0.5),
         "Comparative evaluation across Llama, Ministral, Qwen2, and Granite.",
         size=18, color=MUTED)
add_footer(s, 6, 6)

out = "/home/tio/Documents/Starjob/starjob_intro_methodology.pptx"
prs.save(out)
print(f"Saved: {out}")
